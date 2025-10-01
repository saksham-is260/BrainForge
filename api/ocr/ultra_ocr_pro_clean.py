import os
import re
import json
import cv2
import numpy as np
import pytesseract
from PIL import Image
from datetime import datetime
import hashlib
import sys
import fitz  # PyMuPDF
from pptx import Presentation
import traceback

print("🚀 BRAINFORGE OCR PRO - Smart Text Extraction (Torch-Free Optimized)")
print("=" * 60)

class Config:
    MAX_TEXT_LENGTH = 150000
    CACHE_ENABLED = True
    CACHE_DIR = "ocr_cache"
    MIN_CONTENT_LENGTH = 100
    ENABLE_PREPROCESSING = True

# Global variables for engines - TORCH FREE
TESSERACT_READY = False

try:
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    TESSERACT_READY = True
    print("✅ Tesseract Ready")
except:
    try:
        # Try alternative path
        pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
        TESSERACT_READY = True
        print("✅ Tesseract Ready (Linux)")
    except:
        TESSERACT_READY = False
        print("⚠️ Tesseract Not Found")

# EasyOCR completely removed - TORCH FREE
print("ℹ️ EasyOCR disabled (Torch-Free Mode) - Using Tesseract only")

class UltraOCREngine:
    def __init__(self):
        self.cache_dir = Config.CACHE_DIR
        if Config.CACHE_ENABLED and not os.path.exists(self.cache_dir):
            os.makedirs(self.cache_dir)
        self.stats = {
            'files_processed': 0,
            'total_chars_extracted': 0,
            'success_rate': 0
        }
    
    def extract_structured_content(self, text, source_type):
        """Enhanced content structuring for better AI processing"""
        if not text:
            return text
        
        # Extract different content types based on source
        if source_type == 'pdf':
            return self._structure_pdf_content(text)
        elif source_type == 'ppt':
            return self._structure_ppt_content(text)
        else:
            return self._structure_general_content(text)
    
    def _structure_pdf_content(self, text):
        """Structure PDF content with page-wise extraction"""
        lines = text.split('\n')
        structured_content = []
        current_page = 1
        current_section = ""
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Detect page breaks
            if re.match(r'^\s*page\s*\d+\s*$', line.lower()) or re.match(r'^\s*\d+\s*$', line):
                if current_section:
                    structured_content.append({
                        'page': current_page,
                        'content': current_section.strip(),
                        'type': 'text'
                    })
                    current_section = ""
                current_page += 1
                continue
            
            # Detect headings
            if (len(line) < 100 and 
                (line.isupper() or 
                 re.match(r'^[0-9IVX]+\.', line) or
                 re.match(r'^[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*$', line))):
                if current_section:
                    structured_content.append({
                        'page': current_page,
                        'content': current_section.strip(),
                        'type': 'text'
                    })
                    current_section = ""
                structured_content.append({
                    'page': current_page,
                    'content': line,
                    'type': 'heading'
                })
            else:
                current_section += line + " "
        
        # Add remaining content
        if current_section:
            structured_content.append({
                'page': current_page,
                'content': current_section.strip(),
                'type': 'text'
            })
        
        return self._format_structured_content(structured_content)
    
    def _structure_ppt_content(self, text):
        """Structure PPT content with slide-wise extraction - FIXED slide numbering bug"""
        lines = text.split('\n')
        structured_content = []
        current_slide = 1  # Start from 1, ensure no negative
        current_content = ""
        is_title = True
        slide_detected = False
        
        print(f"🔧 Structuring PPT: Processing {len(lines)} lines")
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Enhanced slide detection
            slide_match = re.match(r'^(?:## |Slide\s*|\-\-\- SLIDE\s*)(\d+)', line, re.IGNORECASE)
            if slide_match:
                detected_slide = int(slide_match.group(1))
                print(f"🔍 Detected slide separator: {line} -> Slide {detected_slide}")
                
                if current_content:
                    structured_content.append({
                        'slide': current_slide,
                        'content': current_content.strip(),
                        'type': 'slide_content'
                    })
                    current_content = ""
                
                # Ensure slide number is positive
                if detected_slide > 0 and detected_slide >= current_slide:
                    current_slide = detected_slide
                else:
                    current_slide += 1
                    print(f"⚠️ Invalid slide number {detected_slide}, incrementing to {current_slide}")
                
                is_title = True
                slide_detected = True
                continue
            
            # Fallback detection
            if not slide_detected and re.search(r'slide\s*\d+', line.lower()):
                if current_content:
                    structured_content.append({
                        'slide': current_slide,
                        'content': current_content.strip(),
                        'type': 'slide_content'
                    })
                    current_content = ""
                current_slide += 1
                print(f"🔍 Fallback slide detection: Starting slide {current_slide}")
                is_title = True
                slide_detected = True
                continue
            
            if is_title and line:
                structured_content.append({
                    'slide': current_slide,
                    'content': line,
                    'type': 'title'
                })
                is_title = False
                slide_detected = True
            elif line:
                current_content += line + " "
        
        # Add remaining content
        if current_content:
            structured_content.append({
                'slide': current_slide,
                'content': current_content.strip(),
                'type': 'slide_content'
            })
        
        print(f"✅ PPT Structured: {len(structured_content)} items, final slide: {current_slide}")
        return self._format_structured_content(structured_content)
    
    def _structure_general_content(self, text):
        """Structure general content with paragraph detection"""
        paragraphs = re.split(r'\n\s*\n', text)
        structured_content = []
        
        for i, para in enumerate(paragraphs):
            para = para.strip()
            if not para or len(para) < 10:
                continue
            
            # Classify paragraph type
            if len(para) < 150 and (para.endswith(':') or para.isupper()):
                para_type = 'heading'
            elif re.match(r'^\d+\.', para) or re.match(r'^[•\-*]', para):
                para_type = 'list'
            else:
                para_type = 'paragraph'
            
            structured_content.append({
                'section': i + 1,
                'content': para,
                'type': para_type
            })
        
        return self._format_structured_content(structured_content)
    
    def _format_structured_content(self, structured_data):
        """Format structured data into readable text"""
        formatted_text = ""
        
        for item in structured_data:
            if 'page' in item:
                formatted_text += f"\n--- PAGE {item['page']} ---\n"
            elif 'slide' in item:
                formatted_text += f"\n--- SLIDE {item['slide']} ---\n"
            
            if item['type'] == 'heading':
                formatted_text += f"HEADING: {item['content']}\n\n"
            elif item['type'] == 'title':
                formatted_text += f"TITLE: {item['content']}\n\n"
            else:
                formatted_text += f"CONTENT: {item['content']}\n\n"
        
        return formatted_text.strip()
    
    def clean_text_for_ai(self, text):
        """Enhanced cleaning specifically for AI processing"""
        if not text:
            return ""
        
        # Remove extraction artifacts but keep structure
        patterns_to_remove = [
            r'TABLE/STRUCTURE:.*?\n',
            r'IMAGE DETECTED:.*?\n',
            r'$$EXTRACTION STATS:.*?$$',
            r'\b\d+\b\s*$',
            r'^\s*\w+\s+\d+\s*$',
            r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\$$\$$,]|(?:%[0-9a-fA-F][0-9a-fA-F]))+',
        ]
        
        for pattern in patterns_to_remove:
            text = re.sub(pattern, '', text, flags=re.MULTILINE)
        
        # Clean up whitespace while preserving structure
        text = re.sub(r' +', ' ', text)
        text = re.sub(r'\n\s+', '\n', text)
        
        # Remove very short lines that are noise
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if (line.startswith('--- ') or 
                line.startswith('HEADING:') or
                line.startswith('TITLE:') or
                line.startswith('CONTENT:') or
                (len(line) > 25 and 
                 not line.isdigit() and 
                 not re.match(r'^[\W\d_]+$', line) and
                 len(re.findall(r'[a-zA-Z]', line)) > 5)):
                cleaned_lines.append(line)
        
        text = '\n'.join(cleaned_lines)
        
        return text
    
    def file_hash(self, file_path):
        stats = os.stat(file_path)
        data = f"{file_path}_{stats.st_size}_{stats.st_mtime}"
        return hashlib.md5(data.encode()).hexdigest()
    
    def enhance_image_quality(self, image):
        """Enhanced image preprocessing"""
        try:
            # Convert to grayscale if needed
            if len(image.shape) == 3:
                gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            else:
                gray = image
            
            # Noise reduction
            denoised = cv2.bilateralFilter(gray, 9, 75, 75)
            
            # Contrast enhancement
            clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
            contrast_enhanced = clahe.apply(denoised)
            
            # Sharpening
            kernel = np.array([[-1,-1,-1], [-1,9,-1], [-1,-1,-1]])
            sharpened = cv2.filter2D(contrast_enhanced, -1, kernel)
            
            return sharpened
        except Exception as e:
            print(f"Image enhancement error: {e}")
            return image
    
    def process_pdf(self, pdf_path):
        """Enhanced PDF extraction with structured content"""
        try:
            doc = fitz.open(pdf_path)
            full_text = ""
            page_contents = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract text with better formatting
                text = page.get_text("text", sort=True)
                
                if text.strip():
                    # Clean page text
                    page_text = text.strip()
                    
                    # Remove common PDF artifacts but keep structure
                    artifacts = [
                        r'http[s]?://\S+',
                        r'www\.\S+',
                        r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
                    ]
                    
                    for artifact in artifacts:
                        page_text = re.sub(artifact, '', page_text)
                    
                    page_contents.append({
                        'page': page_num + 1,
                        'content': page_text
                    })
                    
                    full_text += f"\n--- PAGE {page_num + 1} ---\n{page_text}\n"
            
            doc.close()
            
            # Structure the content
            structured_text = self.extract_structured_content(full_text, 'pdf')
            
            if len(structured_text.strip()) < Config.MIN_CONTENT_LENGTH:
                print("⚠️ PDF text extraction low, trying OCR method...")
                return self._extract_pdf_with_ocr(pdf_path)
            
            return self.clean_text_for_ai(structured_text)
            
        except Exception as e:
            return f"PDF extraction error: {str(e)}"
    
    def _extract_pdf_with_ocr(self, pdf_path):
        """Fallback PDF extraction using OCR"""
        try:
            doc = fitz.open(pdf_path)
            full_text = ""
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                
                # Convert to OpenCV format
                nparr = np.frombuffer(img_data, np.uint8)
                img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                
                if img is not None:
                    # Preprocess image
                    enhanced_img = self.enhance_image_quality(img)
                    
                    # Extract text with Tesseract only (TORCH FREE)
                    text_results = []
                    
                    # Tesseract only - no EasyOCR
                    if TESSERACT_READY:
                        try:
                            tesseract_text = pytesseract.image_to_string(enhanced_img, config='--psm 6')
                            if tesseract_text.strip():
                                text_results.append(tesseract_text)
                        except Exception as e:
                            print(f"Tesseract page {page_num} error: {e}")
                    
                    # Use the best result
                    if text_results:
                        best_text = max(text_results, key=len)
                        if len(best_text.strip()) > 30:
                            full_text += f"\n--- PAGE {page_num + 1} ---\n{best_text.strip()}\n"
                
                # Free memory
                del img, enhanced_img
            
            doc.close()
            
            # Structure the OCR content
            structured_text = self.extract_structured_content(full_text, 'pdf')
            return self.clean_text_for_ai(structured_text)
            
        except Exception as e:
            return f"PDF OCR extraction error: {str(e)}"
    
    def process_ppt(self, ppt_path):
        """Enhanced PPT extraction with structured content"""
        try:
            prs = Presentation(ppt_path)
            slides_content = []
            full_text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_title = ""
                
                print(f"🔍 Processing slide {slide_num + 1}")
                
                # Get title
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
                    if len(title) > 2:
                        slide_title = title
                        slide_text.append(f"TITLE: {title}")
                
                # Get content from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if (len(text) > 10 and 
                            not text.isupper() and
                            not re.match(r'^[\d\W]+$', text) and
                            len(re.findall(r'[a-zA-Z]', text)) > 3):
                            slide_text.append(f"CONTENT: {text}")
                
                # Get notes
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame.text.strip():
                    notes = notes_slide.notes_slide.notes_text_frame.text.strip()
                    if len(notes) > 20:
                        slide_text.append(f"NOTES: {notes}")
                
                if slide_text:
                    slide_content = "\n".join(slide_text)
                    slides_content.append({
                        'slide': slide_num + 1,
                        'title': slide_title,
                        'content': slide_content
                    })
                    full_text += f"\n--- SLIDE {slide_num + 1} ---\n{slide_content}\n"
            
            # Structure the content
            structured_text = self.extract_structured_content(full_text, 'ppt')
            
            if len(structured_text.strip()) < Config.MIN_CONTENT_LENGTH:
                return "Insufficient content extracted from presentation. Please ensure the file contains readable text content."
            
            return self.clean_text_for_ai(structured_text)
            
        except Exception as e:
            return f"PPT extraction error: {str(e)}"
    
    def preprocess_image(self, image_path):
        """Enhanced image preprocessing for better OCR"""
        techniques = {}
        
        try:
            img = cv2.imread(image_path)
            if img is None:
                return techniques
            
            # Basic grayscale
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            techniques['grayscale'] = gray
            
            # Enhanced preprocessing
            enhanced = self.enhance_image_quality(img)
            techniques['enhanced'] = enhanced
            
            # Multiple thresholding techniques
            _, thresh_binary = cv2.threshold(enhanced, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            techniques['threshold_binary'] = thresh_binary
            
            thresh_adaptive = cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                                   cv2.THRESH_BINARY, 11, 2)
            techniques['threshold_adaptive'] = thresh_adaptive
            
            # Morphological operations to clean up text
            kernel = np.ones((2,2), np.uint8)
            morph_open = cv2.morphologyEx(thresh_binary, cv2.MORPH_OPEN, kernel)
            techniques['morphological'] = morph_open
            
        except Exception as e:
            print(f"Image preprocessing error: {e}")
        
        return techniques
    
    def extract_with_tesseract(self, image, config_name, config):
        """Enhanced Tesseract extraction"""
        try:
            if not TESSERACT_READY:
                return None
            
            # Additional config for better accuracy
            enhanced_config = f'{config} -c preserve_interword_spaces=1'
            text = pytesseract.image_to_string(image, config=enhanced_config)
            
            if text.strip():
                # Calculate confidence
                data = pytesseract.image_to_data(image, config=enhanced_config, output_type=pytesseract.Output.DICT)
                confidence = np.mean([float(conf) for conf in data['conf'] if float(conf) > 0]) / 100.0
                
                return {
                    'engine': f'tesseract_{config_name}',
                    'text': text.strip(),
                    'confidence': confidence if not np.isnan(confidence) else 0.7,
                    'length': len(text)
                }
        except Exception as e:
            print(f"Tesseract error ({config_name}): {e}")
        return None
    
    def smart_text_fusion(self, results):
        """Enhanced text fusion from multiple engines"""
        if not results:
            return ""
        
        # Filter out low-quality results
        filtered_results = [r for r in results if r['confidence'] > 0.3 and r['length'] > 10]
        
        if not filtered_results:
            # If all results are low quality, take the best one anyway
            filtered_results = results
        
        # Sort by confidence and length with weights
        filtered_results.sort(key=lambda x: (x['confidence'] * 0.7 + min(x['length']/1000, 0.3)), reverse=True)
        
        # Take the best result
        best_result = filtered_results[0]
        return best_result['text']
    
    def extract_from_image(self, image_path):
        """Enhanced main image extraction - TORCH FREE"""
        print(f"🔍 Processing image: {os.path.basename(image_path)}")
        
        # Cache check
        cache_key = self.file_hash(image_path)
        cache_file = os.path.join(self.cache_dir, f"{cache_key}.json")
        
        if Config.CACHE_ENABLED and os.path.exists(cache_file):
            print("📦 Using cached result")
            with open(cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        
        # Preprocess
        preprocessed = self.preprocess_image(image_path)
        if not preprocessed:
            return {"error": "Image processing failed"}
        
        print(f"🛠️ Using {len(preprocessed)} preprocessing techniques")
        
        # Extract with multiple techniques - TESSERACT ONLY (TORCH FREE)
        all_results = []
        tesseract_configs = [
            ('psm3', '--psm 3'),
            ('psm6', '--psm 6'),
            ('psm4', '--psm 4'),
            ('psm8', '--psm 8'),
        ]
        
        for technique_name, image in preprocessed.items():
            print(f"  🔧 Technique: {technique_name}")
            
            # Tesseract with different configs only
            for config_name, config in tesseract_configs:
                tesseract_result = self.extract_with_tesseract(image, config_name, config)
                if tesseract_result:
                    tesseract_result['technique'] = technique_name
                    all_results.append(tesseract_result)
                    print(f"    ✅ Tesseract-{config_name}: {len(tesseract_result['text'])} chars, conf: {tesseract_result['confidence']:.2f}")
        
        # Combine results
        fused_text = self.smart_text_fusion(all_results)
        
        # Structure the extracted content
        structured_text = self.extract_structured_content(fused_text, 'image')
        cleaned_text = self.clean_text_for_ai(structured_text)
        
        # Update stats
        self.stats['files_processed'] += 1
        self.stats['total_chars_extracted'] += len(cleaned_text)
        self.stats['success_rate'] = (self.stats['success_rate'] * (self.stats['files_processed'] - 1) + 1) / self.stats['files_processed']
        
        result = {
            'text': cleaned_text,
            'source': 'image',
            'techniques_used': len(preprocessed),
            'engines_used': len(all_results),
            'timestamp': datetime.now().isoformat(),
            'stats': {
                'characters': len(cleaned_text),
                'sentences': len(re.split(r'[.!?]+', cleaned_text)),
                'words': len(re.findall(r'\b\w+\b', cleaned_text)),
                'extraction_confidence': max([r['confidence'] for r in all_results]) if all_results else 0,
                'structure_quality': 'high' if '--- PAGE' in cleaned_text or '--- SLIDE' in cleaned_text else 'medium'
            },
            'processing_info': {
                'preprocessing_techniques': list(preprocessed.keys()),
                'engines_tried': list(set([r['engine'] for r in all_results]))
            }
        }
        
        # Cache result
        if Config.CACHE_ENABLED:
            with open(cache_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False)
            print(f"💾 Result cached: {cache_file}")
        
        return result
    
    def process_file(self, file_path):
        """Enhanced main function to process any file type"""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        file_ext = os.path.splitext(file_path)[1].lower()
        file_size = os.path.getsize(file_path) / (1024 * 1024)  # MB
        
        print(f"📁 Processing: {os.path.basename(file_path)} ({file_size:.2f} MB)")
        
        # File size check
        if file_size > 50:  # 50MB limit
            return {"error": "File too large. Please use files smaller than 50MB."}
        
        try:
            if file_ext == '.pdf':
                text = self.process_pdf(file_path)
                source_type = 'pdf'
            elif file_ext in ['.ppt', '.pptx']:
                text = self.process_ppt(file_path)
                source_type = 'ppt'
            elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp']:
                result = self.extract_from_image(file_path)
                if 'error' in result:
                    return result
                text = result['text']
                source_type = 'image'
            else:
                return {"error": f"Unsupported file type: {file_ext}. Supported: PDF, PPT/PPTX, JPG, PNG, BMP, TIFF, WEBP"}
            
            # Enhanced validation
            if not text or len(text.strip()) < Config.MIN_CONTENT_LENGTH:
                return {
                    "error": f"Insufficient text extracted ({len(text.strip()) if text else 0} characters). " +
                            "Please try a different file with clearer text content."
                }
            
            return {
                "success": True,
                "text": text,
                "source_type": source_type,
                "filename": os.path.basename(file_path),
                "processed_at": datetime.now().isoformat(),
                "length": len(text),
                "message": f"Successfully extracted {len(text)} characters from {source_type.upper()}",
                "quality_metrics": {
                    "characters": len(text),
                    "words": len(re.findall(r'\b\w+\b', text)),
                    "sentences": len(re.split(r'[.!?]+', text)),
                    "content_density": len(text) / max(1, len(re.findall(r'\s+', text))),
                    "structure_preserved": 'high' if '--- PAGE' in text or '--- SLIDE' in text else 'medium'
                }
            }
            
        except Exception as e:
            print(f"❌ Processing error: {e}")
            traceback.print_exc()
            return {"error": f"Processing failed: {str(e)}"}
    
    def get_stats(self):
        """Get processing statistics"""
        return self.stats

def main():
    if len(sys.argv) != 2:
        print("Usage: python ultra_ocr.py <filename>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    engine = UltraOCREngine()
    
    result = engine.process_file(file_path)
    
    if "error" in result:
        print(f"❌ Error: {result['error']}")
        sys.exit(1)
    
    print(f"✅ Successfully processed {result['filename']}")
    print(f"📝 Text length: {result['length']} characters")
    print(f"📊 Quality metrics: {result.get('quality_metrics', {})}")
    print(f"📄 Sample structure: {result['text'][:500]}...")

if __name__ == "__main__":
    main()